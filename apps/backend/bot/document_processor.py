import tempfile
from datetime import datetime, timezone
from pathlib import Path

from bot.gemini_ocr import ocr_image_bytes, ocr_pdf_file

_ocr_image_bytes = ocr_image_bytes  # alias for internal callers

SUPPORTED_EXTENSIONS = {".pdf", ".docx", ".pptx"}


def detect_file_type(filename: str) -> str | None:
    ext = Path(filename).suffix.lower()
    if ext == ".pdf":
        return "pdf"
    elif ext == ".docx":
        return "docx"
    elif ext == ".pptx":
        return "pptx"
    return None


def ocr_image_file(file_path: str) -> str:
    with open(file_path, "rb") as f:
        return ocr_image_bytes(f.read())


def convert_to_pdf(input_path: str, output_dir: str) -> str:
    import comtypes.client
    import pythoncom

    pythoncom.CoInitialize()

    src = Path(input_path).resolve()
    out = Path(output_dir).resolve() / f"{src.stem}.pdf"
    ext = src.suffix.lower()

    try:
        if ext == ".docx":
            app = comtypes.client.CreateObject("Word.Application")
            app.Visible = False
            doc = app.Documents.Open(str(src))
            doc.ExportAsFixedFormat(str(out), 17)
            doc.Close()
            app.Quit()

        elif ext == ".pptx":
            app = comtypes.client.CreateObject("PowerPoint.Application")
            app.Visible = False
            pres = app.Presentations.Open(str(src))
            pres.SaveAs(str(out), 32)
            pres.Close()
            app.Quit()

        else:
            raise ValueError(f"Unsupported for COM conversion: {ext}")

    finally:
        pythoncom.CoUninitialize()

    return str(out)


def _extract_pdf(path: str) -> list[dict]:
    try:
        text = ocr_pdf_file(path)
        if text.strip():
            return [{"page": 1, "source": "gemini", "text": text}]
    except Exception:
        pass
    import fitz
    doc = fitz.open(path)
    try:
        pages = []
        for page_num, page in enumerate(doc):
            native_text = page.get_text().strip()
            pix = page.get_pixmap(dpi=200)
            img_bytes = pix.tobytes("png")
            ocr_text = _ocr_image_bytes(img_bytes)
            combined = f"{native_text}\n{ocr_text}".strip()
            source = "native+ocr" if native_text and ocr_text else ("ocr" if ocr_text else "native")
            pages.append({
                "page": page_num + 1,
                "source": source,
                "text": combined,
            })
        return pages
    finally:
        doc.close()


def _try_gemini_pdf(path: str) -> str | None:
    with tempfile.TemporaryDirectory() as tmpdir:
        try:
            pdf_path = convert_to_pdf(path, tmpdir)
            text = ocr_pdf_file(pdf_path)
            if text.strip():
                return text
        except Exception:
            pass
    return None


def _extract_docx(path: str) -> list[dict]:
    text = _try_gemini_pdf(path)
    if text is not None:
        return [{"page": 1, "source": "gemini", "text": text}]
    from docx import Document
    doc = Document(path)
    paragraphs = [p.text for p in doc.paragraphs]
    native_text = "\n".join(paragraphs).strip()
    pages = [{"page": 1, "source": "native", "text": native_text}]
    ocr_parts = []
    for rel in doc.part.rels.values():
        if "image" in rel.reltype:
            image_data = rel.target_part.blob
            text = _ocr_image_bytes(image_data)
            if text.strip():
                ocr_parts.append(text)
    if ocr_parts:
        pages.append({"page": 1, "source": "native+ocr", "text": "\n".join(ocr_parts)})
    return pages


def _extract_pptx(path: str) -> list[dict]:
    text = _try_gemini_pdf(path)
    if text is not None:
        return [{"page": 1, "source": "gemini", "text": text}]
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    prs = Presentation(path)
    slides_data = []
    for slide_num, slide in enumerate(prs.slides):
        texts = []
        for shape in slide.shapes:
            if shape.has_text_frame and shape.text_frame.text.strip():
                texts.append(shape.text_frame.text)
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                img_data = shape.image.blob
                ocr_text = _ocr_image_bytes(img_data)
                if ocr_text.strip():
                    texts.append(ocr_text)
        slides_data.append({
            "page": slide_num + 1,
            "source": "native",
            "text": "\n".join(texts),
        })
    return slides_data


def process_document(
    file_path: str,
    chat_id: int,
    message_id: int,
    sender_name: str,
    sent_at: datetime,
    chat_title: str | None = None,
) -> dict:
    path = Path(file_path)
    file_type = detect_file_type(path.name)

    if file_type == "pdf":
        pages = _extract_pdf(file_path)
    elif file_type == "docx":
        pages = _extract_docx(file_path)
    elif file_type == "pptx":
        pages = _extract_pptx(file_path)
    else:
        raise ValueError(f"Unsupported file type: {file_type}")

    full_text = "\n".join(p["text"] for p in pages).strip()

    return {
        "chat_id": chat_id,
        "message_id": message_id,
        "sender_name": sender_name,
        "sent_at": sent_at.isoformat(),
        "chat_title": chat_title,
        "file_name": path.name,
        "file_type": file_type,
        "processed_at": datetime.now(timezone.utc).isoformat(),
        "pages": pages,
        "full_text": full_text,
    }
